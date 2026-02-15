from __future__ import annotations

import asyncio
import contextlib
import logging
import re
from dataclasses import dataclass
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Callable, Optional

import fitz
import httpx
import numpy as np
import pymupdf4llm
from docx import Document as DocxDocument
from fastembed import TextEmbedding
from openpyxl import load_workbook
from pptx import Presentation
from rank_bm25 import BM25Okapi
from PIL import Image

from .config import rag_config
from storage import get_storage_service

logger = logging.getLogger(__name__)


class DocumentStatus(str, Enum):
    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"


@dataclass
class DocumentProcessingInfo:
    doc_id: str
    filename: str
    status: DocumentStatus = DocumentStatus.PENDING
    progress: int = 0
    progress_message: str = ""
    error: Optional[str] = None
    started_at: Optional[datetime] = None
    completed_at: Optional[datetime] = None
    file_size_bytes: int = 0
    file_extension: str = ""
    file_path: Optional[str] = None


@dataclass
class IndexedDocument:
    doc_id: str
    filename: str
    chunks: list[str]
    bm25: BM25Okapi
    embeddings: np.ndarray
    embedding_norms: np.ndarray


@dataclass
class RetrievedChunk:
    doc_id: str
    filename: str
    index: int
    text: str
    score: float


class RAGService:
    _instance: Optional["RAGService"] = None

    _initialized: bool
    _lock: asyncio.Lock
    _documents_processed: bool
    _document_status: dict[str, DocumentProcessingInfo]
    _document_indexes: dict[str, IndexedDocument]
    _embedder: Optional[TextEmbedding]
    _storage: Optional[Any]

    def __init__(self) -> None:
        if RAGService._instance is not None:
            return
        self._initialized = False
        self._lock = asyncio.Lock()
        self._documents_processed = False
        self._document_status = {}
        self._document_indexes = {}
        self._embedder = None
        self._storage = None
        RAGService._instance = self

    @classmethod
    def get_instance(cls) -> "RAGService":
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance

    async def _get_storage(self):
        if self._storage is None:
            self._storage = await get_storage_service()
        return self._storage

    def register_document(self, doc_id: str, filename: str) -> DocumentProcessingInfo:
        info = DocumentProcessingInfo(
            doc_id=doc_id,
            filename=filename,
            status=DocumentStatus.PENDING,
            started_at=datetime.now(),
        )
        self._document_status[doc_id] = info
        return info

    def register_document_with_metadata(
        self, doc_id: str, filename: str, file_path: str, file_size_bytes: int
    ) -> DocumentProcessingInfo:
        file_extension = Path(filename).suffix.lstrip(".").lower()
        info = DocumentProcessingInfo(
            doc_id=doc_id,
            filename=filename,
            status=DocumentStatus.PENDING,
            started_at=datetime.now(),
            file_path=file_path,
            file_size_bytes=file_size_bytes,
            file_extension=file_extension,
        )
        self._document_status[doc_id] = info
        return info

    def get_document_status(self, doc_id: str) -> Optional[DocumentProcessingInfo]:
        return self._document_status.get(doc_id)

    def update_document_progress(self, doc_id: str, progress: int, message: str = "") -> None:
        info = self._document_status.get(doc_id)
        if info is None:
            return
        info.progress = min(100, max(0, progress))
        info.progress_message = message

    def list_documents(self) -> list[DocumentProcessingInfo]:
        return list(self._document_status.values())

    def _update_document_status(
        self, doc_id: str, status: DocumentStatus, error: Optional[str] = None
    ) -> None:
        info = self._document_status.get(doc_id)
        if info is None:
            return
        info.status = status
        if error:
            info.error = error
        if status in (DocumentStatus.COMPLETED, DocumentStatus.FAILED):
            info.completed_at = datetime.now()

    async def initialize(self) -> None:
        async with self._lock:
            if self._initialized:
                return
            self._embedder = await asyncio.to_thread(
                TextEmbedding, model_name=rag_config.fastembed_model
            )
            self._initialized = True
            logger.info("RAG service initialized with FastEmbed model=%s", rag_config.fastembed_model)

    async def _ensure_initialized(self) -> None:
        if not self._initialized:
            await self.initialize()

    async def process_document(
        self,
        file_path: str,
        doc_id: Optional[str] = None,
        filename: Optional[str] = None,
        on_progress: Optional[Callable[[int, str], None]] = None,
    ) -> dict[str, Any]:
        await self._ensure_initialized()

        def emit_progress(progress: int, message: str) -> None:
            if doc_id:
                self.update_document_progress(doc_id, progress, message)
            if on_progress:
                on_progress(progress, message)

        if doc_id:
            self._update_document_status(doc_id, DocumentStatus.PROCESSING)

        try:
            path = Path(file_path)
            emit_progress(10, "Extracting text")
            text = await asyncio.to_thread(self._extract_text, path)

            emit_progress(40, "Chunking document")
            chunks = self._chunk_text(text)
            if not chunks:
                raise ValueError("No extractable text found in document")

            emit_progress(65, "Building BM25 index")
            tokenized_chunks = [self._tokenize(chunk) for chunk in chunks]
            bm25 = BM25Okapi(tokenized_chunks)

            emit_progress(80, "Building semantic index")
            embeddings = await asyncio.to_thread(self._embed_texts, chunks)
            norms = np.linalg.norm(embeddings, axis=1)
            norms = np.where(norms == 0.0, 1e-12, norms)

            effective_filename = filename or path.name
            indexed_doc = IndexedDocument(
                doc_id=doc_id or path.stem,
                filename=effective_filename,
                chunks=chunks,
                bm25=bm25,
                embeddings=embeddings,
                embedding_norms=norms,
            )

            key = doc_id or path.stem
            self._document_indexes[key] = indexed_doc
            self._documents_processed = True

            await self._persist_parsed_text(key, text)

            if doc_id:
                self._update_document_status(doc_id, DocumentStatus.COMPLETED)
                emit_progress(100, "Processing complete")

            return {"success": True, "file_path": file_path, "doc_id": key}
        except Exception as exc:
            error_msg = str(exc)
            logger.error("Failed to process document %s: %s", file_path, error_msg)
            if doc_id:
                self._update_document_status(doc_id, DocumentStatus.FAILED, error_msg)
            return {"success": False, "error": error_msg, "file_path": file_path}

    async def query(
        self,
        question: str,
        mode: str = "hybrid",
        top_k: int = 10,
        doc_ids: Optional[list[str]] = None,
    ) -> str:
        del mode
        await self._ensure_initialized()

        if not self._document_indexes:
            return ""

        retrieve_top_k = max(1, top_k, rag_config.final_context_k)
        timeout_seconds = max(5, rag_config.query_timeout_seconds)

        try:
            candidates = await asyncio.wait_for(
                asyncio.to_thread(
                    self._retrieve_candidates,
                    question,
                    doc_ids,
                    retrieve_top_k,
                ),
                timeout=timeout_seconds,
            )
        except asyncio.TimeoutError as exc:
            raise TimeoutError("RAG query timed out") from exc

        if not candidates:
            return ""

        context_text = self._build_context_blob(candidates, rag_config.synthesis_max_chars)
        if not rag_config.groq_api_key:
            return context_text

        prompt = (
            "Answer the user question using only the retrieved document excerpts. "
            "If the answer is not in the excerpts, say so clearly.\n\n"
            f"Question: {question}\n\n"
            f"Retrieved excerpts:\n{context_text}"
        )
        synthesized = await self._call_groq(prompt)
        return synthesized if synthesized.strip() else context_text

    async def get_context_for_topic(
        self,
        topic: str,
        mode: str = "hybrid",
        doc_ids: Optional[list[str]] = None,
    ) -> str:
        del mode
        await self._ensure_initialized()

        if not self._document_indexes:
            logger.info("No indexed documents, skipping context retrieval")
            return ""

        retrieve_top_k = max(12, rag_config.final_context_k)
        timeout_seconds = max(5, rag_config.query_timeout_seconds)

        try:
            candidates = await asyncio.wait_for(
                asyncio.to_thread(
                    self._retrieve_candidates,
                    topic,
                    doc_ids,
                    retrieve_top_k,
                ),
                timeout=timeout_seconds,
            )
        except asyncio.TimeoutError as exc:
            raise TimeoutError("RAG context retrieval timed out") from exc

        if not candidates:
            return ""

        context_text = self._build_context_blob(candidates, rag_config.synthesis_max_chars)
        if not rag_config.groq_api_key:
            return context_text

        prompt = (
            "Create concise presentation-ready context from retrieved excerpts. "
            "Focus on facts, key concepts, useful examples, and concrete data.\n\n"
            f"Topic: {topic}\n\n"
            f"Retrieved excerpts:\n{context_text}"
        )
        synthesized = await self._call_groq(prompt)
        return synthesized if synthesized.strip() else context_text

    async def delete_document(self, doc_id: str) -> bool:
        info = self._document_status.get(doc_id)
        if info is None:
            return False

        storage = await self._get_storage()

        if info.file_path:
            filename = Path(info.file_path).name
            await storage.delete_upload(filename)

        await storage.delete_parsed(doc_id)

        self._document_indexes.pop(doc_id, None)
        del self._document_status[doc_id]
        self._documents_processed = bool(self._document_indexes)
        return True

    def get_config_info(self) -> dict[str, Any]:
        return {
            "working_dir": str(rag_config.working_dir),
            "parser": rag_config.parser,
            "embedding_model": rag_config.fastembed_model,
            "embedding_dim": rag_config.embedding_dim,
            "llm_model": rag_config.groq_model if rag_config.groq_api_key else "disabled",
            "initialized": self._initialized,
            "has_documents": bool(self._document_indexes),
        }

    def _extract_text(self, path: Path) -> str:
        ext = path.suffix.lower()
        if ext == ".pdf":
            return self._extract_pdf_text(path)
        if ext == ".docx":
            return self._extract_docx_text(path)
        if ext == ".pptx":
            return self._extract_pptx_text(path)
        if ext == ".xlsx":
            return self._extract_xlsx_text(path)
        if ext in {".txt", ".md"}:
            return path.read_text(encoding="utf-8", errors="ignore")
        if ext in {".png", ".jpg", ".jpeg"}:
            return self._extract_image_metadata(path)
        if ext in {".doc", ".ppt", ".xls"}:
            return self._extract_legacy_binary_text(path, ext)
        raise ValueError(f"Unsupported file extension: {ext}")

    def _extract_pdf_text(self, path: Path) -> str:
        markdown = pymupdf4llm.to_markdown(str(path))
        cleaned = self._sanitize_text(markdown)
        if cleaned:
            return cleaned

        with fitz.open(str(path)) as doc:
            page_text = [page.get_text("text") for page in doc]
        fallback = self._sanitize_text("\n\n".join(page_text))
        if not fallback:
            raise ValueError("No extractable text found in PDF")
        return fallback

    def _extract_docx_text(self, path: Path) -> str:
        doc = DocxDocument(str(path))
        parts: list[str] = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_values:
                    parts.append(" | ".join(row_values))
        combined = self._sanitize_text("\n\n".join(parts))
        if not combined:
            raise ValueError("No extractable text found in DOCX")
        return combined

    def _extract_pptx_text(self, path: Path) -> str:
        deck = Presentation(str(path))
        slide_text: list[str] = []
        for slide_index, slide in enumerate(deck.slides, start=1):
            collected: list[str] = [f"Slide {slide_index}"]
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    collected.append(text.strip())
            if len(collected) > 1:
                slide_text.append("\n".join(collected))
        combined = self._sanitize_text("\n\n".join(slide_text))
        if not combined:
            raise ValueError("No extractable text found in PPTX")
        return combined

    def _extract_xlsx_text(self, path: Path) -> str:
        workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
        entries: list[str] = []
        for sheet in workbook.worksheets:
            entries.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    entries.append(" | ".join(values))
        workbook.close()
        combined = self._sanitize_text("\n".join(entries))
        if not combined:
            raise ValueError("No extractable text found in XLSX")
        return combined

    def _extract_image_metadata(self, path: Path) -> str:
        with Image.open(path) as img:
            width, height = img.size
            mode = img.mode
            image_format = img.format or path.suffix.lstrip(".").upper()
        return (
            f"Image file: {path.name}\n"
            f"Format: {image_format}\n"
            f"Resolution: {width}x{height}\n"
            f"Color mode: {mode}\n"
            "OCR is not enabled in this fast pipeline."
        )

    def _extract_legacy_binary_text(self, path: Path, ext: str) -> str:
        raw = path.read_bytes()
        decoded = raw.decode("latin-1", errors="ignore")
        printable = "".join(ch if ch.isprintable() else " " for ch in decoded)
        words = re.findall(r"[A-Za-z0-9][A-Za-z0-9_\-/.]{2,}", printable)
        if len(words) >= 200:
            return self._sanitize_text(" ".join(words[:5000]))
        return (
            f"Legacy binary office file {path.name} ({ext}) was uploaded. "
            "Convert to modern format (.docx, .pptx, .xlsx) for higher retrieval quality."
        )

    def _sanitize_text(self, text: str) -> str:
        normalized = text.replace("\r\n", "\n").replace("\r", "\n")
        normalized = re.sub(r"\n{3,}", "\n\n", normalized)
        normalized = re.sub(r"[ \t]{2,}", " ", normalized)
        return normalized.strip()

    def _chunk_text(self, text: str) -> list[str]:
        clean = self._sanitize_text(text)
        if not clean:
            return []

        paragraphs = [block.strip() for block in re.split(r"\n{2,}", clean) if block.strip()]
        expanded: list[str] = []
        for paragraph in paragraphs:
            if len(paragraph) <= rag_config.chunk_size_chars:
                expanded.append(paragraph)
                continue
            expanded.extend(self._split_large_paragraph(paragraph, rag_config.chunk_size_chars))

        chunks: list[str] = []
        current = ""
        for paragraph in expanded:
            candidate = paragraph if not current else f"{current}\n\n{paragraph}"
            if len(candidate) <= rag_config.chunk_size_chars:
                current = candidate
                continue

            if current and len(current) >= rag_config.min_chunk_chars:
                chunks.append(current)

            overlap = current[-rag_config.chunk_overlap_chars :] if current else ""
            overlap = overlap.strip()
            current = f"{overlap}\n\n{paragraph}".strip() if overlap else paragraph
            if len(current) > rag_config.chunk_size_chars:
                chunks.append(current[: rag_config.chunk_size_chars].strip())
                current = current[rag_config.chunk_size_chars - rag_config.chunk_overlap_chars :].strip()

        if current and len(current) >= rag_config.min_chunk_chars:
            chunks.append(current)

        if not chunks:
            chunks = [clean[: rag_config.chunk_size_chars]]
        return chunks

    def _split_large_paragraph(self, paragraph: str, max_chars: int) -> list[str]:
        sentences = [part.strip() for part in re.split(r"(?<=[.!?])\s+", paragraph) if part.strip()]
        if not sentences:
            return [paragraph]

        pieces: list[str] = []
        current = ""
        for sentence in sentences:
            candidate = sentence if not current else f"{current} {sentence}"
            if len(candidate) <= max_chars:
                current = candidate
                continue
            if current:
                pieces.append(current)
            current = sentence
        if current:
            pieces.append(current)
        return pieces

    def _tokenize(self, text: str) -> list[str]:
        tokens = re.findall(r"[a-z0-9]+", text.lower())
        return tokens if tokens else ["_empty_"]

    def _embed_texts(self, texts: list[str]) -> np.ndarray:
        if self._embedder is None:
            raise RuntimeError("Embedder is not initialized")
        vectors = list(self._embedder.embed(texts))
        if not vectors:
            raise RuntimeError("Embedding generation returned no vectors")
        array = np.vstack([np.asarray(vector, dtype=np.float32) for vector in vectors])
        return array

    def _embed_query(self, query: str) -> np.ndarray:
        if self._embedder is None:
            raise RuntimeError("Embedder is not initialized")
        vector = list(self._embedder.embed([query]))
        if not vector:
            raise RuntimeError("Embedding generation failed for query")
        return np.asarray(vector[0], dtype=np.float32)

    def _select_documents(self, doc_ids: Optional[list[str]]) -> list[IndexedDocument]:
        if not doc_ids:
            return list(self._document_indexes.values())

        selected: list[IndexedDocument] = []
        seen: set[str] = set()
        for doc_id in doc_ids:
            if doc_id in seen:
                continue
            seen.add(doc_id)
            indexed = self._document_indexes.get(doc_id)
            if indexed is not None:
                selected.append(indexed)
        return selected

    def _retrieve_candidates(
        self,
        query: str,
        doc_ids: Optional[list[str]],
        top_k: int,
    ) -> list[RetrievedChunk]:
        selected_docs = self._select_documents(doc_ids)
        if not selected_docs:
            return []

        query_tokens = self._tokenize(query)
        query_embedding = self._embed_query(query)
        query_norm = float(np.linalg.norm(query_embedding))
        if query_norm == 0.0:
            query_norm = 1e-12

        all_candidates: list[RetrievedChunk] = []
        per_retriever_k = max(top_k, rag_config.per_retriever_k)

        for doc in selected_docs:
            if not doc.chunks:
                continue

            bm25_scores = np.asarray(doc.bm25.get_scores(query_tokens), dtype=np.float32)
            dense_raw = np.dot(doc.embeddings, query_embedding)
            dense_scores = dense_raw / (doc.embedding_norms * query_norm)

            candidate_count = min(len(doc.chunks), per_retriever_k)
            if candidate_count <= 0:
                continue

            bm25_order = np.argsort(-bm25_scores)[:candidate_count]
            dense_order = np.argsort(-dense_scores)[:candidate_count]

            merged_scores: dict[int, float] = {}
            for rank, chunk_index in enumerate(bm25_order, start=1):
                merged_scores[int(chunk_index)] = merged_scores.get(int(chunk_index), 0.0) + (
                    1.0 / (rag_config.rrf_k + rank)
                )
            for rank, chunk_index in enumerate(dense_order, start=1):
                merged_scores[int(chunk_index)] = merged_scores.get(int(chunk_index), 0.0) + (
                    1.0 / (rag_config.rrf_k + rank)
                )

            ranked = sorted(merged_scores.items(), key=lambda item: item[1], reverse=True)
            for chunk_index, score in ranked[:candidate_count]:
                all_candidates.append(
                    RetrievedChunk(
                        doc_id=doc.doc_id,
                        filename=doc.filename,
                        index=chunk_index,
                        text=doc.chunks[chunk_index],
                        score=score,
                    )
                )

        all_candidates.sort(key=lambda candidate: candidate.score, reverse=True)
        return all_candidates[:top_k]

    def _build_context_blob(self, candidates: list[RetrievedChunk], max_chars: int) -> str:
        blocks: list[str] = []
        for position, candidate in enumerate(candidates, start=1):
            blocks.append(
                f"[{position}] Source: {candidate.filename} (doc_id={candidate.doc_id}, chunk={candidate.index})\n"
                f"{candidate.text}"
            )
        merged = "\n\n---\n\n".join(blocks)
        if len(merged) <= max_chars:
            return merged
        return merged[:max_chars]

    async def _call_groq(self, prompt: str) -> str:
        if not rag_config.groq_api_key:
            return ""

        payload = {
            "model": rag_config.groq_model,
            "temperature": 0.1,
            "max_tokens": rag_config.groq_max_tokens,
            "messages": [
                {
                    "role": "system",
                    "content": "You are a retrieval-grounded assistant. Use only provided excerpts.",
                },
                {"role": "user", "content": prompt},
            ],
        }
        headers = {
            "Authorization": f"Bearer {rag_config.groq_api_key}",
            "Content-Type": "application/json",
        }

        try:
            timeout = httpx.Timeout(rag_config.llm_timeout_seconds)
            async with httpx.AsyncClient(timeout=timeout) as client:
                response = await client.post(
                    rag_config.groq_chat_completions_url,
                    headers=headers,
                    json=payload,
                )
                response.raise_for_status()
                data = response.json()
            return (
                data.get("choices", [{}])[0]
                .get("message", {})
                .get("content", "")
                .strip()
            )
        except Exception as exc:
            logger.warning("Groq synthesis failed: %s", exc)
            return ""

    async def _persist_parsed_text(self, doc_id: str, text: str) -> None:
        storage = await self._get_storage()
        await storage.save_parsed(doc_id, text)


rag_service = RAGService()


async def get_rag_service() -> RAGService:
    await rag_service.initialize()
    return rag_service
