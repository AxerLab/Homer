"""
PPTX to Images converter using python-pptx and Pillow
"""
import os
import json
from pathlib import Path
from typing import List, Dict, Any
import tempfile
import shutil
from pptx import Presentation
from PIL import Image
import io
import base64

def extract_pptx_slides_info(pptx_path: str) -> Dict[str, Any]:
    """
    Extract slide information from a PPTX file.
    Returns slide titles, text content, and basic metadata.
    """
    try:
        prs = Presentation(pptx_path)
        slides_info = []

        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "title": "",
                "content": [],
                "shapes_count": len(slide.shapes)
            }

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        # First non-empty text is likely the title
                        if not slide_data["title"] and len(text) < 200:
                            slide_data["title"] = text
                        else:
                            slide_data["content"].append(text)

            # If no title found, use slide number
            if not slide_data["title"]:
                slide_data["title"] = f"Slide {i + 1}"

            slides_info.append(slide_data)

        return {
            "total_slides": len(prs.slides),
            "slides": slides_info,
            "width": prs.slide_width,
            "height": prs.slide_height
        }
    except Exception as e:
        print(f"Error extracting PPTX info: {e}")
        return {
            "error": str(e),
            "total_slides": 0,
            "slides": []
        }

def convert_pptx_to_html_preview(pptx_path: str) -> Dict[str, Any]:
    """
    Convert PPTX to a simplified HTML preview by extracting text content.
    This provides a basic preview without requiring image conversion.
    """
    info = extract_pptx_slides_info(pptx_path)

    if "error" in info:
        return info

    # Create HTML representation of slides
    html_slides = []
    for slide in info["slides"]:
        slide_html = f"""
        <div class="slide-preview" data-slide="{slide['slide_number']}">
            <div class="slide-header">
                <span class="slide-number">Slide {slide['slide_number']}</span>
            </div>
            <div class="slide-content">
                <h2 class="slide-title">{slide['title']}</h2>
                <div class="slide-body">
                    {''.join([f'<p>{content}</p>' for content in slide['content']])}
                </div>
            </div>
        </div>
        """
        html_slides.append(slide_html)

    return {
        "total_slides": info["total_slides"],
        "slides": info["slides"],
        "html_preview": ''.join(html_slides),
        "preview_type": "text"
    }

def get_pptx_preview(file_path: str) -> Dict[str, Any]:
    """
    Main function to get PPTX preview data.
    Returns structured data that can be rendered in the frontend.
    """
    if not os.path.exists(file_path):
        return {"error": "File not found"}

    # For now, use HTML text preview
    # In the future, this could be extended to generate actual slide images
    return convert_pptx_to_html_preview(file_path)